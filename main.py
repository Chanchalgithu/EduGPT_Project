import streamlit as st
import pickle, faiss, os, random, datetime, time
import numpy as np
import pandas as pd
from openai import OpenAI
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer

# -------- File Handling --------
import fitz  # PyMuPDF for PDF
import docx
from pptx import Presentation
from PIL import Image

# ---------- SETTINGS ----------
load_dotenv()
API_KEY = os.getenv("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY")
MODEL_LLM = "gpt-4o-mini"
INDEX_PATH = "embeddings/embeddings_index.faiss"
TEXTS_PATH = "embeddings/embeddings_texts.pkl"
EMBED_MODEL = "all-MiniLM-L6-v2"
TOP_K = 3
USERS_FILE = "users.csv"
SESSION_FILE = "session.csv"
SESSION_DURATION = 6 * 60 * 60   # 6 hours

# ---------- CACHE LOADERS ----------
@st.cache_resource
def load_index_texts():
    index = faiss.read_index(INDEX_PATH)
    texts = pickle.load(open(TEXTS_PATH, "rb"))
    embedder = SentenceTransformer(EMBED_MODEL)
    return index, texts, embedder

index, texts, embedder = load_index_texts()
client = OpenAI(api_key=API_KEY)

# ---------- FILE PARSER ----------
def parse_file(uploaded_file):
    try:
        ext = uploaded_file.name.split(".")[-1].lower()

        if ext == "pdf":
            text = ""
            doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            for page in doc:
                text += page.get_text()
            return text

        elif ext in ["docx", "doc"]:
            doc = docx.Document(uploaded_file)
            return "\n".join([p.text for p in doc.paragraphs])

        elif ext in ["xlsx", "xls", "csv"]:
            df = pd.read_excel(uploaded_file) if ext != "csv" else pd.read_csv(uploaded_file)
            return df.to_string()

        elif ext in ["pptx", "ppt"]:
            prs = Presentation(uploaded_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        elif ext in ["png", "jpg", "jpeg"]:
            img = Image.open(uploaded_file)
            return f"[Image Uploaded: {uploaded_file.name}, Size={img.size}]"

        elif ext in ["mp4", "avi", "mov", "mkv"]:
            return f"[Video Uploaded: {uploaded_file.name}]"

        else:
            return f"[Unsupported file type: {ext}]"

    except Exception as e:
        return f"[Error reading file: {e}]"

# ---------- QUERY HANDLER ----------
def answer_query(query, extra_context=""):
    try:
        q_emb = embedder.encode([query])
        D, I = index.search(np.array(q_emb).astype("float32"), TOP_K)
        retrieved = [texts[i] for i in I[0] if i < len(texts)]
        dataset_context = "\n".join(retrieved).strip()

        context = (dataset_context + "\n\n" + extra_context).strip()

        if context:
            prompt = f"Use the following dataset and file context to answer:\n\n{context}\n\nQuestion: {query}\n\nAnswer:"
        else:
            prompt = f"Answer this question:\n\n{query}"

        resp = client.chat.completions.create(
            model=MODEL_LLM,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=500
        )
        return resp.choices[0].message.content.strip()

    except Exception as e:
        return f"⚠️ Error: {e}"

# ---------- USER MANAGEMENT ----------
def load_users():
    if os.path.exists(USERS_FILE):
        return pd.read_csv(USERS_FILE)
    else:
        return pd.DataFrame(columns=["username", "password", "phone"])

def save_user(username, password, phone):
    df = load_users()
    if username in df["username"].values:
        return False
    new_row = pd.DataFrame([[username, password, phone]], columns=df.columns)
    df = pd.concat([df, new_row], ignore_index=True)
    df.to_csv(USERS_FILE, index=False)
    return True

# ---------- SESSION MANAGEMENT ----------
def load_session():
    if os.path.exists(SESSION_FILE):
        df = pd.read_csv(SESSION_FILE)
        if len(df) > 0:
            return {"username": df.iloc[0]["username"], "login_time": float(df.iloc[0]["login_time"])}
    return None

def save_session(username):
    now = time.time()
    df = pd.DataFrame([[username, now]], columns=["username", "login_time"])
    df.to_csv(SESSION_FILE, index=False)

def clear_session():
    if os.path.exists(SESSION_FILE):
        os.remove(SESSION_FILE)

def is_session_valid():
    session = load_session()
    if session:
        now = time.time()
        if now - session["login_time"] < SESSION_DURATION:
            st.session_state["logged_in"] = True
            st.session_state["username"] = session["username"]
            return True
        else:
            clear_session()
    return False

# ---------- OTP ----------
def send_otp(phone):
    otp = str(random.randint(1000, 9999))
    st.session_state["otp_store"] = {phone: otp}
    return otp

# ---------- LOGIN / SIGNUP ----------
def login_signup():
    st.markdown("<h2 style='text-align: center;'>🔑 EduGPT Login / Signup</h2>", unsafe_allow_html=True)
    choice = st.radio("Choose option:", ["Login", "Signup"], horizontal=True)

    if choice == "Signup":
        with st.form("signup_form", clear_on_submit=False):
            username = st.text_input("Create Username")
            password = st.text_input("Create Password", type="password")
            phone = st.text_input("Phone Number")

            send_otp_btn = st.form_submit_button("📩 Send OTP")
            if send_otp_btn:
                if phone:
                    otp = send_otp(phone)
                    st.success(f"OTP sent to {phone} (demo OTP: {otp})")
                else:
                    st.error("Enter phone number first!")

            otp_in = st.text_input("Enter OTP")
            signup_btn = st.form_submit_button("✅ Signup Now")

            if signup_btn:
                if "otp_store" in st.session_state and phone in st.session_state["otp_store"]:
                    if st.session_state["otp_store"][phone] == otp_in:
                        if save_user(username, password, phone):
                            st.success("✅ Account created! Please login.")
                        else:
                            st.error("User already exists.")
                    else:
                        st.error("❌ Invalid OTP!")
                else:
                    st.error("Please request OTP first.")

    elif choice == "Login":
        with st.form("login_form", clear_on_submit=False):
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            login_btn = st.form_submit_button("🔓 Login")

            if login_btn:
                df = load_users()
                if ((df["username"] == username) & (df["password"] == password)).any():
                    st.session_state["logged_in"] = True
                    st.session_state["username"] = username
                    save_session(username)  # ✅ save login session
                    st.success(f"🎉 Welcome {username} (session valid for 6 hours)")
                else:
                    st.error("Invalid credentials")

# ---------- MAIN APP ----------
def main_app():
    st.set_page_config(page_title="EduGPT ⚡", layout="wide")

    # ✅ Responsive CSS for mobile/tablet
    st.markdown("""
        <style>
            @media (max-width: 768px) {
                .block-container {
                    padding: 1rem !important;
                }
                .stTextInput, .stTextArea, .stButton, .stFileUploader {
                    width: 100% !important;
                }
            }
        </style>
    """, unsafe_allow_html=True)

    # ---- Session state for chat history ----
    if "history" not in st.session_state:
        st.session_state.history = {}

    today = datetime.date.today().strftime("%Y-%m-%d")
    if today not in st.session_state.history:
        st.session_state.history[today] = []

    # ---- SIDEBAR ----
    with st.sidebar:
        st.header("💬 Chat History")
        for day, chats in st.session_state.history.items():
            with st.expander(day, expanded=(day == today)):
                for i, chat in enumerate(chats):
                    st.markdown(f"**Q{i+1}:** {chat['q']}")
                    if st.button(f"🗑 Delete", key=f"del_{day}_{i}"):
                        st.session_state.history[day].pop(i)
                        st.rerun()

        if st.button("🆕 New Chat"):
            st.session_state.history[today] = []
            st.rerun()

        if st.button("🗑 Clear All"):
            st.session_state.history = {today: []}
            st.rerun()

        if st.button("🚪 Logout"):
            clear_session()
            st.session_state["logged_in"] = False
            st.experimental_rerun()

    # ---- Main Page ----
    st.title("⚡ EduGPT – Super Fast Mode")

    uploaded_file = st.file_uploader("📎 Upload File (PDF, Excel, PPT, DOCX, Image, Video)", type=None)

    file_context = ""
    if uploaded_file:
        st.success(f"✅ File Uploaded: {uploaded_file.name}")
        file_context = parse_file(uploaded_file)

    with st.form("chat_form", clear_on_submit=True):
        query = st.text_input("❓ Ask your question:", key="query_input")
        submitted = st.form_submit_button("Get Answer (Press Enter)")

    if submitted and query.strip():
        with st.spinner("⚡ Thinking..."):
            ans = answer_query(query, file_context)

        st.session_state.history[today].append({"q": query, "a": ans})
        st.markdown("### ✅ Answer:")
        st.write(ans)

# ---------- APP FLOW ----------
if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False

# ✅ auto-login if session valid
if not st.session_state["logged_in"]:
    if is_session_valid():
        st.session_state["logged_in"] = True

if st.session_state["logged_in"]:
    main_app()
else:
    login_signup()


